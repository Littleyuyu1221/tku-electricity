from __future__ import annotations

import argparse
import ast
import json
import re
import subprocess
import sys
import zipfile
from dataclasses import asdict, dataclass
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


SCRIPT_DIR = Path(__file__).resolve().parent
if (SCRIPT_DIR / "app.py").is_file():
    ROOT = SCRIPT_DIR
    OUTPUTS = SCRIPT_DIR
else:
    ROOT = SCRIPT_DIR.parent
    OUTPUTS = ROOT / "outputs"


@dataclass
class Check:
    status: str
    item: str
    evidence: str


def add(checks: list[Check], condition: bool, item: str, evidence: str, failure: str) -> None:
    checks.append(Check("PASS" if condition else "FAIL", item, evidence if condition else failure))


def latest_package() -> Path | None:
    packages = list(OUTPUTS.glob("tku-electricity-專題完整包-v*.zip"))
    if not packages:
        return None

    def version_number(path: Path) -> int:
        match = re.search(r"-v(\d+)\.zip$", path.name)
        return int(match.group(1)) if match else -1

    return max(packages, key=version_number)


def run_checks() -> list[Check]:
    checks: list[Check] = []
    required = [
        "app.py",
        "README.md",
        "requirements.txt",
        "大四專題報告.docx",
        "大四專題口試簡報.pptx",
        "口試講稿與答辯手冊.md",
        "宿舍用電資料範本.csv",
        "宿舍用電與冷氣資料蒐集範本.xlsx",
        "資料蒐集與驗收清單.md",
        "校方資料申請與研究使用說明.docx",
        "校方資料申請信範本.md",
        "專題送審完成度報告.md",
        "預覽-首頁.png",
        "預覽-用電預測.png",
        "預覽-冷氣風險.png",
    ]
    missing = [name for name in required if not (OUTPUTS / name).is_file()]
    add(checks, not missing, "核心交付物", f"{len(required)} 項必要檔案皆存在", "缺少：" + "、".join(missing))

    app_path = OUTPUTS / "app.py"
    app_text = app_path.read_text(encoding="utf-8")
    try:
        ast.parse(app_text, filename=str(app_path))
        syntax_ok = True
        syntax_error = ""
    except SyntaxError as exc:
        syntax_ok = False
        syntax_error = f"第 {exc.lineno} 行：{exc.msg}"
    add(checks, syntax_ok, "程式語法", "app.py 可由 Python 完整解析", syntax_error)

    page_names = [
        "① 快速看懂用電",
        "② 住宿人數與開館",
        "③ 電都用到哪裡",
        "④ 可以省多少電",
        "⑤ 下個月用電預測",
        "⑥ 冷氣機型與故障風險",
    ]
    missing_pages = [name for name in page_names if name not in app_text]
    add(checks, not missing_pages, "六個展示頁面", "六個功能頁面皆在程式中註冊", "未找到：" + "、".join(missing_pages))
    add(checks, "高齡" not in app_text, "需求範圍", "程式中已無高齡專用功能或文字", "仍找到『高齡』文字")

    model_validator = ROOT / "validate_models.py"
    if not model_validator.is_file():
        model_validator = ROOT / "work" / "validate_models.py"
    model_result = subprocess.run(
        [sys.executable, str(model_validator)],
        cwd=ROOT,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        timeout=60,
        check=False,
    )
    model_output = (model_result.stdout + "\n" + model_result.stderr).strip()
    add(
        checks,
        model_result.returncode == 0 and "MODEL_VALIDATION_OK" in model_output,
        "模型回歸測試",
        "用電預測、冷氣耗電與故障風險測試通過",
        model_output[-500:] or "模型驗證沒有輸出",
    )

    report = Document(OUTPUTS / "大四專題報告.docx")
    report_text = "\n".join(p.text for p in report.paragraphs)
    report_headings = [p.text for p in report.paragraphs if p.style.name.startswith("Heading")]
    report_keywords = ["摘要", "研究方法", "示範資料結果", "研究限制", "參考文獻", "模型卡"]
    missing_report = [key for key in report_keywords if key not in report_text]
    report_ok = (
        len(report.paragraphs) >= 150
        and len(report_headings) >= 40
        and len(report.tables) >= 10
        and len(report.inline_shapes) >= 4
        and not missing_report
    )
    add(
        checks,
        report_ok,
        "正式專題報告結構",
        f"{len(report.paragraphs)} 段、{len(report_headings)} 個標題、{len(report.tables)} 表、{len(report.inline_shapes)} 圖",
        "報告結構不足或缺少章節：" + "、".join(missing_report),
    )

    deck = Presentation(OUTPUTS / "大四專題口試簡報.pptx")
    note_texts = [
        slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
        for slide in deck.slides
    ]
    deck_ok = len(deck.slides) == 12 and all(note_texts) and all("[Sources]" in text for text in note_texts)
    add(
        checks,
        deck_ok,
        "口試簡報",
        f"{len(deck.slides)} 頁；{sum(bool(text) for text in note_texts)} 頁有講者備註與來源",
        "簡報頁數、講者備註或來源標記不完整",
    )

    workbook = load_workbook(OUTPUTS / "宿舍用電與冷氣資料蒐集範本.xlsx", data_only=False)
    expected_sheets = {"使用說明", "月用電資料", "冷氣設備清冊", "保養故障紀錄", "資料字典", "品質檢核", "選單設定"}
    formulas = 0
    validations = 0
    for sheet in workbook.worksheets:
        formulas += sum(
            1
            for row in sheet.iter_rows()
            for cell in row
            if isinstance(cell.value, str) and cell.value.startswith("=")
        )
        validations += len(sheet.data_validations.dataValidation)
    workbook_ok = expected_sheets.issubset(set(workbook.sheetnames)) and formulas > 0 and validations > 0
    add(
        checks,
        workbook_ok,
        "校方資料蒐集工作簿",
        f"{len(workbook.sheetnames)} 個工作表、{formulas} 個公式、{validations} 組輸入驗證",
        "工作表、公式或輸入驗證不完整",
    )

    package = latest_package()
    if package is None and OUTPUTS == ROOT:
        checks.append(Check("PASS", "完整交付包", f"目前從已解壓的完整交付包執行，{len(required) + 2} 項檔案由其他檢查逐項驗證"))
    elif package is None:
        checks.append(Check("FAIL", "完整交付包", "找不到版本化 ZIP"))
    else:
        with zipfile.ZipFile(package) as archive:
            names = set(archive.namelist())
        archive_required = required + ["final_readiness.py", "validate_models.py"]
        archive_missing = [name for name in archive_required if name not in names]
        add(
            checks,
            not archive_missing,
            "完整交付包",
            f"{package.name} 含 {len(names)} 個交付檔案",
            "壓縮包缺少：" + "、".join(archive_missing),
        )

    report_has_identity_placeholder = "［請填寫系所名稱］" in report_text or "［姓名／學號］" in report_text
    deck_text = "\n".join(
        shape.text
        for slide in deck.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    deck_has_identity_placeholder = "＿＿＿＿＿＿" in deck_text
    checks.append(
        Check(
            "PENDING" if report_has_identity_placeholder or deck_has_identity_placeholder else "PASS",
            "學生與指導資料",
            "報告或簡報仍有姓名、學號、系所、指導教授待填欄位"
            if report_has_identity_placeholder or deck_has_identity_placeholder
            else "身分欄位已填妥",
        )
    )
    checks.append(Check("PENDING", "校方正式資料", "目前仍為示範資料；須取得授權的用電、住宿、冷氣清冊與維修故障紀錄後重跑模型"))
    checks.append(Check("PENDING", "校系格式與指導老師確認", "送件前仍需套用系所指定封面／頁碼格式並取得指導老師內容確認"))
    checks.append(Check("WARN", "Word 逐頁視覺檢查", "此電腦缺少 Word／LibreOffice 轉頁工具；已完成結構、表格與無障礙檢查，但未完成逐頁影像驗證"))
    return checks


def print_markdown(checks: list[Check]) -> None:
    print("| 狀態 | 檢查項目 | 證據／待辦 |")
    print("|---|---|---|")
    for check in checks:
        evidence = check.evidence.replace("|", "／").replace("\n", " ")
        print(f"| {check.status} | {check.item} | {evidence} |")


def main() -> int:
    parser = argparse.ArgumentParser(description="驗證大四專題本機交付物完整度")
    parser.add_argument("--json", action="store_true", help="以 JSON 輸出")
    args = parser.parse_args()
    checks = run_checks()
    if args.json:
        print(json.dumps([asdict(check) for check in checks], ensure_ascii=False, indent=2))
    else:
        print_markdown(checks)
    failed = [check for check in checks if check.status == "FAIL"]
    print("\nLOCAL_READINESS_OK" if not failed else "\nLOCAL_READINESS_FAILED")
    return 1 if failed else 0


if __name__ == "__main__":
    raise SystemExit(main())
